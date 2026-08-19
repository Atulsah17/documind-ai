"""Verify multi-format extraction (files are generated in-memory)."""
import io
import json

from app.rag.loaders import extract_text


def test_txt_and_md():
    assert "hello" in extract_text("a.txt", b"hello world")
    assert "Title" in extract_text("a.md", b"# Title\n\ncontent")


def test_csv_and_tsv():
    out = extract_text("data.csv", b"name,role\nAtul,Engineer")
    assert "Atul" in out and "Engineer" in out
    out2 = extract_text("data.tsv", b"name\trole\nAtul\tEngineer")
    assert "Atul" in out2


def test_json():
    out = extract_text("d.json", json.dumps({"framework": "fastapi"}).encode())
    assert "fastapi" in out


def test_html():
    html = b"<html><body><h1>Heading</h1><script>ignore()</script><p>Body text</p></body></html>"
    out = extract_text("page.html", html)
    assert "Heading" in out and "Body text" in out and "ignore" not in out


def test_docx():
    import docx

    d = docx.Document()
    d.add_paragraph("Kubernetes orchestrates containers.")
    buf = io.BytesIO()
    d.save(buf)
    out = extract_text("doc.docx", buf.getvalue())
    assert "Kubernetes" in out


def test_xlsx():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["metric", "value"])
    ws.append(["accuracy", 0.95])
    buf = io.BytesIO()
    wb.save(buf)
    out = extract_text("sheet.xlsx", buf.getvalue())
    assert "accuracy" in out and "0.95" in out


def test_pptx():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Review"
    buf = io.BytesIO()
    prs.save(buf)
    out = extract_text("deck.pptx", buf.getvalue())
    assert "Quarterly Review" in out


def test_unknown_extension_falls_back_to_text():
    assert "print" in extract_text("script.py", b"print('hi')")
    assert "data" in extract_text("weird.xyz", b"just some data")
