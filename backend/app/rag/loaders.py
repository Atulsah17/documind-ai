"""Multi-format document text extraction.

Dispatches by file extension to the right parser and always degrades gracefully
to a best-effort text decode so no upload hard-fails.
"""
from __future__ import annotations

import csv
import io
import json

# Extensions we can parse well (advertised to the UI's file picker).
SUPPORTED_EXTENSIONS = [
    ".pdf", ".docx", ".pptx", ".xlsx", ".xls",
    ".csv", ".tsv", ".json", ".html", ".htm",
    ".md", ".markdown", ".txt", ".rst", ".log",
    # images & scanned docs (read via vision model)
    ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp",
    # common code / config files
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rb",
    ".c", ".cpp", ".cs", ".sh", ".yaml", ".yml", ".toml", ".ini", ".sql",
]


def _decode(raw: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="ignore")


def _pdf(raw: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(raw: bytes) -> str:
    import docx  # python-docx

    doc = docx.Document(io.BytesIO(raw))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pptx(raw: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def _xlsx(raw: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _csv_like(raw: bytes, delimiter: str) -> str:
    text = _decode(raw)
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    return "\n".join(" | ".join(row) for row in reader if any(cell.strip() for cell in row))


def _html(raw: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(raw), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _json(raw: bytes) -> str:
    try:
        return json.dumps(json.loads(_decode(raw)), indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return _decode(raw)


def extract_text(filename: str, raw: bytes) -> str:
    """Return plain text for any supported document type."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    try:
        if ext == ".pdf":
            return _pdf(raw)
        if ext == ".docx":
            return _docx(raw)
        if ext == ".pptx":
            return _pptx(raw)
        if ext in (".xlsx", ".xls"):
            return _xlsx(raw)
        if ext == ".csv":
            return _csv_like(raw, ",")
        if ext == ".tsv":
            return _csv_like(raw, "\t")
        if ext in (".html", ".htm"):
            return _html(raw)
        if ext == ".json":
            return _json(raw)
    except Exception as exc:  # parser failed → fall back to raw decode
        fallback = _decode(raw)
        if fallback.strip():
            return fallback
        raise ValueError(f"Could not parse {filename}: {exc}") from exc

    # markdown, txt, code, unknown → decode as text
    return _decode(raw)
